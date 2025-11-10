"""
PowerPoint Report Generator
Generates weekly reports in PowerPoint format while preserving formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from typing import Dict, Any
import os


class WeeklyReportGenerator:
    """Generates VLines weekly reports in PowerPoint format"""

    def __init__(self, config: Dict[str, Any]):
        self.config = config
        self.template_path = config.get('template', {}).get('path', './templates/weekly_report_template.pptx')
        self.output_dir = config.get('output', {}).get('directory', './reports')
        self.filename_pattern = config.get('output', {}).get('filename_pattern', 'VLines_Weekly_Report_{date}.pptx')

    def generate_report(self, data: Dict[str, Any]) -> str:
        """
        Generate PowerPoint report from data

        Args:
            data: Dictionary containing all report data

        Returns:
            Path to generated PowerPoint file
        """
        try:
            # Load template
            if os.path.exists(self.template_path):
                prs = Presentation(self.template_path)
                print(f"Loaded template: {self.template_path}")
            else:
                print(f"Template not found: {self.template_path}")
                print("Creating blank presentation (create template for production use)")
                prs = Presentation()
                self._create_sample_slides(prs, data)

            # Update slides with data
            self._update_tckt_slide(prs, data.get('tckt', {}))
            self._update_ops_slide(prs, data.get('ops', {}))
            self._update_kinh_doanh_slide(prs, data.get('kinh_doanh', {}))

            # Generate output filename
            output_path = self._get_output_path(data)

            # Ensure output directory exists
            os.makedirs(self.output_dir, exist_ok=True)

            # Save presentation
            prs.save(output_path)
            print(f"✓ Report generated successfully: {output_path}")

            return output_path

        except Exception as e:
            print(f"✗ Error generating report: {e}")
            raise

    def _update_tckt_slide(self, prs: Presentation, tckt_data: Dict[str, Any]):
        """Update TCKT (Accounting) slide with data"""
        print("Updating TCKT slide...")

        # If you have a template, you would find shapes by name and update them
        # Example approach for template-based updates:
        # for slide in prs.slides:
        #     for shape in slide.shapes:
        #         if shape.has_text_frame:
        #             self._replace_text_in_runs(shape, '{{TOTAL_RECEIVABLES}}',
        #                                       f"{tckt_data.get('overview', {}).get('receivables', {}).get('total', 0):,}")

        # For now, just log the data that would be updated
        overview = tckt_data.get('overview', {})
        if overview:
            receivables = overview.get('receivables', {})
            print(f"  - Total Receivables: {receivables.get('total', 0):,}")
            print(f"  - Within Term: {receivables.get('within_term', 0):,}")
            print(f"  - Overdue: {receivables.get('overdue', 0):,}")

    def _update_ops_slide(self, prs: Presentation, ops_data: Dict[str, Any]):
        """Update OPS (Operations) slide with ship schedule data"""
        print("Updating OPS slide...")

        ship_schedule = ops_data.get('ship_schedule', [])
        print(f"  - Processing {len(ship_schedule)} ships")

        # Template-based update would go here
        # You would find the table shape and update cells

    def _update_kinh_doanh_slide(self, prs: Presentation, kinh_doanh_data: Dict[str, Any]):
        """Update Kinh Doanh (Business) slide with market data"""
        print("Updating Kinh Doanh slide...")

        market_overview = kinh_doanh_data.get('market_overview', {})
        if market_overview:
            hph_hcm = market_overview.get('hph_hcm_route', {})
            print(f"  - HPH-HCM VLines share: {hph_hcm.get('vlines_share', 0)}%")

    def _replace_text_in_runs(self, shape, old_text: str, new_text: str):
        """
        Replace text while preserving formatting (at run level)

        This is the KEY method for preserving fonts and formatting
        """
        if not shape.has_text_frame:
            return

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

    def _create_sample_slides(self, prs: Presentation, data: Dict[str, Any]):
        """Create sample slides when no template is available"""
        print("Creating sample slides (template recommended for production)")

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "VLines Weekly Report"
        week = data.get('metadata', {}).get('week', 'N/A')
        subtitle.text = f"Week {week}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"

        # Content slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = "Report Summary"

        tf = body_shape.text_frame
        tf.text = "This is a sample report"

        p = tf.add_paragraph()
        p.text = "Create a template file (weekly_report_template.pptx) for production use"
        p.level = 1

    def _get_output_path(self, data: Dict[str, Any]) -> str:
        """Generate output file path"""
        date_str = datetime.now().strftime('%Y-%m-%d')
        week = data.get('metadata', {}).get('week', 'unknown')

        filename = self.filename_pattern.format(
            date=date_str,
            week=week
        )

        return os.path.join(self.output_dir, filename)

    def _format_currency(self, value: int) -> str:
        """Format number as Vietnamese currency"""
        return f"{value:,} VNĐ"


if __name__ == "__main__":
    # Test the generator
    from data_fetcher import DataFetcher

    test_config = {
        'data_source': {'type': 'json'},
        'template': {'path': './templates/weekly_report_template.pptx'},
        'output': {
            'directory': './reports',
            'filename_pattern': 'VLines_Weekly_Report_{date}.pptx'
        }
    }

    # Fetch sample data
    fetcher = DataFetcher(test_config)
    data = fetcher.fetch_data()

    # Generate report
    generator = WeeklyReportGenerator(test_config)
    output_path = generator.generate_report(data)

    print(f"\nTest report generated at: {output_path}")
